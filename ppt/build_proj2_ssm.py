from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
TEMPLATE = ROOT / "ppt_template_proj2.pptx"
OUTPUT = ROOT / "proj2_ssm.pptx"

BLUE = RGBColor(0x1C, 0x73, 0xFF)
CYAN = RGBColor(0x50, 0xC4, 0xEE)
NAVY = RGBColor(0x12, 0x2C, 0x59)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x59, 0x59, 0x59)
MID_GRAY = RGBColor(0x98, 0xA2, 0xB3)
LIGHT_GRAY = RGBColor(0xF4, 0xF7, 0xFC)
BAR_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT_TITLE = "Gwangyang Sunshine Bold"
FONT_BODY = "Pretendard Regular"
FONT_BOLD = "Pretendard SemiBold"


def remove_all_slides(prs):
    for slide_id in list(prs.slides._sldIdLst):
        prs.part.drop_rel(slide_id.rId)
        prs.slides._sldIdLst.remove(slide_id)


def no_line(shape):
    shape.line.fill.background()


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def add_shape(slide, kind, x, y, w, h, fill, line=None, radius_adjust=None):
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shape, fill)
    if line is None:
        no_line(shape)
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.5)
    if radius_adjust is not None and shape.adjustments:
        shape.adjustments[0] = radius_adjust
    return shape


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    size=22,
    color=BLACK,
    font=FONT_BODY,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.MIDDLE,
    margin=0.04,
    line_spacing=1.0,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    lines = text.split("\n")
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_rich_title(slide, y, segments):
    box = slide.shapes.add_textbox(Inches(1.6), Inches(y), Inches(23.45), Inches(1.0))
    tf = box.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    for text, color in segments:
        run = p.add_run()
        run.text = text
        run.font.name = FONT_TITLE
        run.font.size = Pt(58)
        run.font.color.rgb = color
    return box


def add_picture_contain(slide, path, x, y, w, h):
    with Image.open(path) as image:
        iw, ih = image.size
    scale = min(w / iw, h / ih)
    pw = iw * scale
    ph = ih * scale
    px = x + (w - pw) / 2
    py = y + (h - ph) / 2
    return slide.shapes.add_picture(str(path), Inches(px), Inches(py), Inches(pw), Inches(ph))


def add_base(slide, section, deck_name="제형 유사도 검색"):
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 26.667, 15, BLUE)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.72, -0.45, 25.23, 14.72, WHITE)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.68, 0.98, 23.72, 0.86, BAR_GRAY)
    add_text(slide, 2.3, 1.08, 7.0, 0.55, section, 19, BLACK)
    add_text(slide, 18.5, 1.08, 6.25, 0.55, deck_name, 18, BLACK, align=PP_ALIGN.RIGHT)


def add_arrow(slide, x, y, w=0.55, h=0.32, color=BLUE):
    return add_shape(slide, MSO_SHAPE.DOWN_ARROW, x, y, w, h, color)


def add_notes(slide, text):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = text


def build_slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_base(slide, "01. 형상 검색 구조")
    add_rich_title(slide, 2.22, [("형상으로 제형 후보를 좁히는 ", BLACK), ("검색 구조", BLUE)])
    add_text(
        slide,
        4.0,
        3.20,
        18.67,
        0.62,
        "알약의 외곽 윤곽을 수치화해 새 이미지와 형태가 가까운 제형을 찾습니다.",
        23,
        DARK_GRAY,
        align=PP_ALIGN.CENTER,
    )

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.65, 4.15, 9.55, 7.72, LIGHT_GRAY)
    add_text(slide, 2.15, 4.42, 3.8, 0.52, "검색 과정", 25, BLACK, FONT_BOLD, True)
    steps = [
        ("01", "알약 이미지 입력"),
        ("02", "알약 영역 분리"),
        ("03", "56차원 형상 임베딩"),
        ("04", "품목 프로토타입 거리 비교"),
        ("05", "29개 제형 중 Top-k 반환"),
    ]
    y = 5.15
    for idx, (num, label) in enumerate(steps):
        fill = WHITE if idx not in (2, 4) else BLUE
        text_color = BLACK if fill == WHITE else WHITE
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 2.12, y, 8.62, 0.92, fill)
        add_shape(slide, MSO_SHAPE.OVAL, 2.38, y + 0.16, 0.60, 0.60, CYAN if idx < 4 else BLACK)
        add_text(slide, 2.38, y + 0.16, 0.60, 0.60, num, 15, WHITE, FONT_BOLD, True, PP_ALIGN.CENTER)
        add_text(slide, 3.22, y + 0.10, 7.10, 0.70, label, 22, text_color, FONT_BOLD, True)
        if idx < len(steps) - 1:
            add_arrow(slide, 6.23, y + 0.98, 0.48, 0.28, BLUE)
        y += 1.27

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 11.63, 4.15, 13.38, 7.72, LIGHT_GRAY)
    add_text(slide, 12.14, 4.42, 5.2, 0.52, "형상 추출 예시", 25, BLACK, FONT_BOLD, True)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 12.12, 5.12, 12.40, 5.85, WHITE)
    add_picture_contain(slide, ROOT / "shape_extraction_process.png", 12.30, 5.28, 12.04, 5.50)
    add_text(
        slide,
        14.68,
        11.05,
        7.20,
        0.56,
        "입력 이미지  →  외곽 윤곽 표현",
        22,
        BLUE,
        FONT_BOLD,
        True,
        PP_ALIGN.CENTER,
    )

    chip_y = 12.25
    chips = [
        (1.65, 7.55, "딥러닝 학습 없이 윤곽 특징 계산"),
        (9.56, 7.55, "한 이미지의 여러 알약 객체 처리"),
        (17.47, 7.55, "회전·반전 영향을 줄여 형상 비교"),
    ]
    for x, w, label in chips:
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, chip_y, w, 1.07, WHITE, BLUE)
        add_text(slide, x + 0.18, chip_y + 0.10, w - 0.36, 0.82, label, 18, BLACK, FONT_BOLD, True, PP_ALIGN.CENTER)

    add_notes(
        slide,
        "저희가 구현한 방식은 알약의 외곽 윤곽을 수치화해 유사한 제형을 검색하는 구조입니다. 오른쪽 예시처럼 입력 이미지에서 알약 영역을 분리하고, 외곽 윤곽을 형상 정보로 변환합니다. 이 윤곽은 56차원 형상 임베딩으로 표현되며, 미리 구축한 품목별 대표 벡터와의 거리를 비교해 29개 제형 가운데 가까운 후보를 순서대로 제시합니다. 딥러닝 모델을 별도로 학습하지 않아도 적용할 수 있고, 알약이 회전되거나 뒤집힌 경우의 차이도 줄이도록 특징을 정규화했습니다. 또한 한 이미지에 앞면과 뒷면이 함께 있는 경우에도 각각의 알약을 찾아 처리할 수 있습니다.",
    )


def build_slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_base(slide, "02. 형상 인덱스")
    add_rich_title(slide, 2.22, [("56차원 특징에서 ", BLACK), ("29개 제형 인덱스", BLUE), ("로", BLACK)])
    add_text(
        slide,
        3.1,
        3.20,
        20.47,
        0.62,
        "윤곽 특징을 결합하고 이미지 정보를 품목과 제형 단위의 대표 벡터로 집계했습니다.",
        23,
        DARK_GRAY,
        align=PP_ALIGN.CENTER,
    )

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.65, 4.18, 14.0, 7.70, LIGHT_GRAY)
    add_text(slide, 2.15, 4.45, 7.2, 0.55, "56차원 형상 임베딩 구성", 25, BLACK, FONT_BOLD, True)
    groups = [
        ("기본 기하 특징", "9차원", "40%", "비율 · 원형도 · 이심률 · 볼록성 · 대칭도", BLUE, 1.0),
        ("EFD", "24차원", "35%", "전체 윤곽의 굴곡과 형태", CYAN, 0.875),
        ("방사형 FFT", "16차원", "20%", "중심에서 외곽까지의 반복 패턴", NAVY, 0.50),
        ("Hu moment", "7차원", "5%", "회전 등에 강한 전역 형상", MID_GRAY, 0.125),
    ]
    y = 5.30
    for name, dim, weight, desc, color, frac in groups:
        add_text(slide, 2.18, y, 3.0, 0.55, name, 21, BLACK, FONT_BOLD, True)
        add_text(slide, 5.13, y, 1.55, 0.55, dim, 18, DARK_GRAY, FONT_BODY)
        add_text(slide, 6.75, y, 4.75, 0.55, desc, 17, DARK_GRAY, FONT_BODY)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 11.55, y + 0.08, 2.75, 0.38, WHITE)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 11.55, y + 0.08, max(0.35, 2.75 * frac), 0.38, color)
        add_text(slide, 14.38, y - 0.02, 0.83, 0.58, weight, 18, color, FONT_BOLD, True, PP_ALIGN.RIGHT)
        y += 1.48

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 16.08, 4.18, 8.93, 7.70, LIGHT_GRAY)
    add_text(slide, 16.58, 4.45, 5.2, 0.55, "인덱스 생성 과정", 25, BLACK, FONT_BOLD, True)
    metrics = [
        ("54,282장", "이미지 특징", CYAN),
        ("27,345개", "품목 프로토타입", BLUE),
        ("29개", "제형 중심 프로토타입", BLACK),
    ]
    y = 5.25
    for idx, (value, label, color) in enumerate(metrics):
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 17.05, y, 6.98, 1.52, WHITE)
        add_text(slide, 17.45, y + 0.15, 3.15, 0.75, value, 31, color, FONT_BOLD, True)
        add_text(slide, 20.35, y + 0.24, 3.20, 0.62, label, 19, DARK_GRAY, FONT_BODY, False, PP_ALIGN.RIGHT)
        if idx < len(metrics) - 1:
            add_arrow(slide, 20.28, y + 1.69, 0.55, 0.34, BLUE)
            add_text(slide, 21.00, y + 1.62, 2.3, 0.46, "중앙값 집계", 15, DARK_GRAY, FONT_BODY)
        y += 2.30

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.65, 12.25, 23.36, 1.10, BLUE)
    add_text(
        slide,
        2.10,
        12.36,
        22.46,
        0.84,
        "제형 중심 + 품목별 대표 벡터를 함께 저장해 다양한 세부 형태를 검색에 반영",
        21,
        WHITE,
        FONT_BOLD,
        True,
        PP_ALIGN.CENTER,
    )

    # PowerPoint export can occasionally place inherited template elements over
    # the header on this dense slide. Re-adding the header last keeps it on top.
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.68, 0.98, 23.72, 0.86, BAR_GRAY)
    add_text(slide, 2.3, 1.08, 7.0, 0.55, "02. 형상 인덱스", 19, BLACK)
    add_text(slide, 21.55, 1.08, 3.15, 0.55, "제형 유사도 검색", 17, BLACK, align=PP_ALIGN.CENTER)

    add_notes(
        slide,
        "형상 임베딩은 네 종류의 특징을 결합한 56차원 벡터입니다. 기본 기하 특징은 가로세로비, 원형도, 이심률, 볼록성, 대칭도처럼 직관적인 모양을 표현합니다. EFD는 윤곽 전체의 굴곡을, 방사형 FFT는 중심에서 외곽까지 거리의 반복 패턴을 담고, Hu moment는 전역적인 형태를 보완합니다. 이렇게 추출한 5만 4,282장의 이미지 특징을 먼저 품목별 중앙값으로 모아 2만 7,345개의 품목 프로토타입을 만들었습니다. 이어서 이를 제형별로 다시 집계해 29개의 제형 중심을 구성했습니다. 검색할 때는 제형 중심뿐 아니라 품목별 대표 벡터도 함께 활용해, 같은 제형 안에서 나타나는 여러 세부 형태까지 반영합니다.",
    )


def build_slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_base(slide, "03. 검색 결과")
    add_rich_title(slide, 2.22, [("실제 이미지에서 확인한 제형 ", BLACK), ("검색 결과", BLUE)])
    add_text(
        slide,
        3.0,
        3.18,
        20.67,
        0.60,
        "서로 다른 외곽 형태의 테스트 이미지에서 기록된 제형이 검색 결과 1위로 나타났습니다.",
        23,
        DARK_GRAY,
        align=PP_ALIGN.CENTER,
    )
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 3.10, 3.88, 20.47, 0.72, LIGHT_GRAY)
    add_text(
        slide,
        3.35,
        3.98,
        19.97,
        0.49,
        "새 이미지 입력  →  객체별 특징 추출  →  품목 프로토타입 거리 계산  →  제형별 Top-k 순위",
        18,
        BLUE,
        FONT_BOLD,
        True,
        PP_ALIGN.CENTER,
    )

    cards = [
        (1.65, "장방형", ROOT / "shape_rectangular_bg.png", BLUE),
        (9.62, "사과형", ROOT / "shape_apple_bg.png", CYAN),
        (17.59, "땅콩형", ROOT / "shape_peanut_bg.png", BLUE),
    ]
    for idx, (x, label, image_path, color) in enumerate(cards, start=1):
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 4.95, 7.42, 6.40, LIGHT_GRAY)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.45, 5.27, 6.52, 0.82, color)
        add_text(slide, x + 0.58, 5.35, 6.26, 0.61, label, 25, WHITE, FONT_BOLD, True, PP_ALIGN.CENTER)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.45, 6.30, 6.52, 3.44, WHITE)
        add_picture_contain(slide, image_path, x + 0.61, 6.44, 6.20, 3.16)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x + 1.47, 9.98, 4.48, 0.78, WHITE, color)
        add_text(slide, x + 1.57, 10.06, 4.28, 0.58, "✓  Top-1 일치", 21, color, FONT_BOLD, True, PP_ALIGN.CENTER)

    info = [
        ("01", "가장 가까운 품목까지의 거리"),
        ("02", "1위와 2위 후보의 거리 차이"),
        ("03", "제형 중심 거리와 유사 품목"),
    ]
    for idx, (num, label) in enumerate(info):
        x = 1.65 + idx * 7.97
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 11.78, 7.42, 1.52, WHITE, BLUE)
        add_shape(slide, MSO_SHAPE.OVAL, x + 0.28, 12.16, 0.76, 0.76, BLUE)
        add_text(slide, x + 0.28, 12.16, 0.76, 0.76, num, 16, WHITE, FONT_BOLD, True, PP_ALIGN.CENTER)
        add_text(slide, x + 1.20, 11.95, 5.88, 1.15, label, 18, BLACK, FONT_BOLD, True)

    add_notes(
        slide,
        "새 이미지가 들어오면 검출된 알약별로 형상 특징을 계산하고, 각 품목 프로토타입까지의 거리를 구한 뒤 제형별 후보 순위를 생성합니다. 화면에는 서로 다른 외곽 형태를 가진 장방형, 사과형, 땅콩형 테스트 이미지를 배치했습니다. 현재 확인한 세 이미지에서는 CSV에 기록된 제형이 각각 검색 결과 1위로 나타났습니다. 결과에는 제형명만 제공되는 것이 아니라 가장 가까운 거리, 1위와 2위의 거리 차이, 제형 중심까지의 거리, 그리고 가장 유사한 실제 품목도 함께 표시됩니다. 따라서 29개 전체 제형을 직접 살펴보는 대신, 형태가 가까운 후보군부터 빠르게 확인할 수 있습니다.",
    )


def main():
    prs = Presentation(str(TEMPLATE))
    remove_all_slides(prs)
    build_slide_1(prs)
    build_slide_2(prs)
    build_slide_3(prs)
    prs.core_properties.title = "알약 제형 유사도 검색"
    prs.core_properties.subject = "형상 임베딩 및 프로토타입 거리 기반 제형 검색"
    prs.core_properties.author = "프로젝트 2팀"
    prs.save(str(OUTPUT))
    print(OUTPUT)


if __name__ == "__main__":
    main()
